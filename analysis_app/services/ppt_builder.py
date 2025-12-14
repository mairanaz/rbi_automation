from __future__ import annotations

import shutil
from dataclasses import dataclass
from pathlib import Path
from typing import Dict, List, Optional, Tuple

from django.conf import settings
from openpyxl import load_workbook
from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.shapes.picture import Picture
from pptx.util import Pt


MASTERFILE_SHEET_NAME = "Masterfile"
FIRST_DATA_ROW = 8


COL_NO = 1
COL_EQUIPMENT_NO = 2
COL_PMT_NO = 3
COL_DESCRIPTION = 4
COL_PARTS = 5
COL_PHASE = 6
COL_FLUID = 7
COL_TYPE = 8
COL_SPEC = 9
COL_GRADE = 10
COL_INSULATION = 11
COL_DESIGN_TEMP = 12
COL_DESIGN_PRESS = 13
COL_OPER_TEMP = 14
COL_OPER_PRESS = 15



TEMPLATES_DIR = Path(settings.BASE_DIR) / "rbi_templates"
INSPECTION_TEMPLATE_FILENAME = "Inspection Plan Template.pptx"
INSPECTION_TEMPLATE_PATH = TEMPLATES_DIR / INSPECTION_TEMPLATE_FILENAME


_EQUIPMENT_SLIDE_MAP: Dict[Tuple[str, str], int] = {
    ("MLK PMT 10101", "V-001"): 0,
    ("MLK PMT 10102", "V-002"): 1,
    ("MLK PMT 10103", "V-003"): 2,
    ("MLK PMT 10104", "V-004"): 3,
    ("MLK PMT 10105", "V-005"): 4,
    ("MLK PMT 10106", "V-006"): 5,
    ("MLK PMT 10107", "H-001"): 6,
    ("MLK PMT 10108", "H-002"): 7,
    ("MLK PMT 10109", "H-003"): 8,
    ("MLK PMT 10110", "H-004"): 9,
}


GENERAL_DESC_BOX = (2914650, 495040, 2514600, 246221)
GENERAL_TAG_BOX  = (5676900, 496864,  990600, 245110)
GENERAL_PMT_BOX  = (7391400, 457200, 1264920, 245110)



@dataclass
class MasterfileRow:
    parts: str
    fluid: str
    type_text: str
    spec: str
    grade: str
    insulation: str
    op_temp: Optional[float]
    op_press: Optional[float]


@dataclass
class EquipmentData:
    description: str
    tag_no: str
    pmt_no: str
    rows: List[MasterfileRow]



def _apply_text_style(tf, text: str) -> None:
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

  
    for p in tf.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        if not p.runs:
            run = p.add_run()
            run.text = text
        for run in p.runs:
            run.font.name = "Arial"
            run.font.size = Pt(8)


def _set_cell_text(cell, text: Optional[str]) -> None:
    
    s = "" if text is None else str(text)
    cell.text = s
    tf = cell.text_frame
    _apply_text_style(tf, s)


def _set_shape_text(shape, text: Optional[str]) -> None:
   
    s = "" if text is None else str(text)
    if not getattr(shape, "has_text_frame", False):
        return
    shape.text_frame.text = s
    _apply_text_style(shape.text_frame, s)


def _norm_pmt(value: str) -> str:
    return " ".join(str(value or "").upper().split())


def _norm_eq(value: str) -> str:
    return str(value or "").upper().replace(" ", "")


def get_template_slide_index(pmt_no: str, equipment_no: str) -> Optional[int]:
    p_norm = _norm_pmt(pmt_no)
    e_norm = _norm_eq(equipment_no)
    for (pmt, eq), idx in _EQUIPMENT_SLIDE_MAP.items():
        if _norm_pmt(pmt) == p_norm and _norm_eq(eq) == e_norm:
            return idx
    return None


def _parse_filename(original_filename: str) -> Tuple[str, str]:
 
    stem = Path(original_filename).stem
    if " - " in stem:
        pmt, eq = stem.split(" - ", 1)
        return pmt.strip(), eq.strip()
    tokens = stem.split()
    if not tokens:
        return stem, ""
    eq = tokens[-1]
    pmt = " ".join(tokens[:-1])
    return pmt.strip(), eq.strip()


def _format_number(value: Optional[float]) -> str:
    if value is None or value == "":
        return ""
    try:
        v = float(value)
    except (TypeError, ValueError):
        return str(value)
    if abs(v - round(v)) < 1e-6:
        return str(int(round(v)))
    return f"{v:.3f}".rstrip("0").rstrip(".")


def _short_type(type_text: Optional[str]) -> str:

    if not type_text:
        return ""
    t = str(type_text).strip().upper()
    if t in {"CS", "SS"}:
        return t
    if "CARBON" in t and "STEEL" in t:
        return "CS"
    if "STAINLESS" in t and "STEEL" in t:
        return "SS"
    return t


def _insulation_to_flag(insulation: Optional[str]) -> str:

    if insulation is None:
        return ""
    text = str(insulation).strip().upper()
    if text in {"Y", "YES"}:
        return "Y"
    if text in {"N", "NO"}:
        return "N"
    if not text:
        return ""
    return "Y"


def _norm_label(s: str) -> str:
    s = str(s or "").upper().strip()
    s = s.replace("\n", " ")
    s = " ".join(s.split())
    return s



def _find_latest_block_start(ws, eq_no: str, pmt_no: str) -> Optional[int]:
    eq_norm = str(eq_no).strip().upper()
    pmt_norm = str(pmt_no).strip().upper()
    max_row = ws.max_row
    start_row: Optional[int] = None

    for r in range(FIRST_DATA_ROW, max_row + 1):
        eq_val = ws.cell(r, COL_EQUIPMENT_NO).value
        pmt_val = ws.cell(r, COL_PMT_NO).value
        if eq_val and pmt_val:
            if str(eq_val).strip().upper() == eq_norm and str(pmt_val).strip().upper() == pmt_norm:
                start_row = r

    return start_row


def _collect_block_rows(ws, start_row: int, eq_no: str, pmt_no: str) -> List[int]:
    rows: List[int] = []
    if not start_row:
        return rows

    max_row = ws.max_row
    target_eq = str(eq_no).strip()
    target_pmt = str(pmt_no).strip()

    r = start_row
    last_with_parts = start_row

    while r <= max_row:
        eq_val = ws.cell(r, COL_EQUIPMENT_NO).value
        pmt_val = ws.cell(r, COL_PMT_NO).value
        parts_val = ws.cell(r, COL_PARTS).value

      
        if r > start_row:
            if (eq_val not in (None, "", target_eq) or pmt_val not in (None, "", target_pmt)):
                break

        
        any_value = any(ws.cell(r, c).value not in (None, "") for c in range(1, COL_OPER_PRESS + 1))
        if not any_value:
            break

        if parts_val not in (None, "", "-"):
            last_with_parts = r

        r += 1

    for rr in range(start_row, last_with_parts + 1):
        rows.append(rr)
    return rows


def _load_equipment_data_from_masterfile(
    masterfile_path: Path,
    eq_no: str,
    pmt_no: str,
) -> EquipmentData:
    wb = load_workbook(masterfile_path, data_only=True)
    if MASTERFILE_SHEET_NAME not in wb.sheetnames:
        raise ValueError(f"Sheet '{MASTERFILE_SHEET_NAME}' not found in {masterfile_path}")
    ws = wb[MASTERFILE_SHEET_NAME]

    start_row = _find_latest_block_start(ws, eq_no, pmt_no)
    if not start_row:
        raise ValueError(f"No rows found for equipment {eq_no} / {pmt_no}")

    block_rows = _collect_block_rows(ws, start_row, eq_no, pmt_no)

    description = str(ws.cell(start_row, COL_DESCRIPTION).value or "").strip()
    rows: List[MasterfileRow] = []

    for r in block_rows:
        rows.append(
            MasterfileRow(
                parts=str(ws.cell(r, COL_PARTS).value or "").strip(),
                fluid=str(ws.cell(r, COL_FLUID).value or "").strip(),
                type_text=str(ws.cell(r, COL_TYPE).value or "").strip(),
                spec=str(ws.cell(r, COL_SPEC).value or "").strip(),
                grade=str(ws.cell(r, COL_GRADE).value or "").strip(),
                insulation=str(ws.cell(r, COL_INSULATION).value or "").strip(),
                op_temp=ws.cell(r, COL_OPER_TEMP).value,
                op_press=ws.cell(r, COL_OPER_PRESS).value,
            )
        )

    return EquipmentData(
        description=description,
        tag_no=str(eq_no),
        pmt_no=str(pmt_no),
        rows=rows,
    )


def _find_material_table(slide):
   
    for shape in slide.shapes:
        if not getattr(shape, "has_table", False):
            continue
        table = shape.table
        if len(table.rows) < 3 or len(table.columns) < 9:
            continue

        tl = table.cell(0, 0).text.strip().upper()
        c1 = table.cell(0, 1).text.strip().upper()
        hdr = " ".join(table.cell(1, 3).text.upper().split())

        if "FLUID" in tl and "COMPONENT" in c1 and "TYPE" in hdr:
            return table
    return None


def _pick_row_by_component(component_text: str, equipment_rows: List[MasterfileRow]) -> Optional[MasterfileRow]:
    comp = _norm_label(component_text)

    for r in equipment_rows:
        p = _norm_label(r.parts)
        if p == comp:
            return r

    for r in equipment_rows:
        p = _norm_label(r.parts)
        if p and comp and (p in comp or comp in p):
            return r

    synonyms = {
        "TOP HEAD": ["HEAD", "TOPHEAD", "DISHED END", "DISHEDEND"],
        "BOTTOM HEAD": ["HEAD", "BOTTOMHEAD", "DISHED END", "DISHEDEND"],
        "HEAD": ["TOP HEAD", "BOTTOM HEAD", "DISHED END", "DISHEDEND"],
        "CHANNEL": ["HEAD", "CHANNEL HEAD", "CHANNELHEAD"],
        "TUBE BUNDLE": ["TUBE", "BUNDLE", "TUBEBUNDLE"],
    }

    for key, syns in synonyms.items():
        if comp == key or any(s in comp for s in syns):
            for r in equipment_rows:
                p = _norm_label(r.parts)
                if p == key or any(s in p for s in syns):
                    return r

    return None


def _fill_material_table(table, data: EquipmentData) -> None:

    if table is None:
        return

    for row_idx in range(2, len(table.rows)):
        cells = table.rows[row_idx].cells

        component_text = cells[1].text  
        src = _pick_row_by_component(component_text, data.rows)

        if not src:
            for col in (0, 3, 4, 5, 6, 7, 8):
                _set_cell_text(cells[col], "")
            continue

        _set_cell_text(cells[0], src.fluid)
        _set_cell_text(cells[3], _short_type(src.type_text))
        _set_cell_text(cells[4], src.spec)
        _set_cell_text(cells[5], src.grade)
        _set_cell_text(cells[6], _insulation_to_flag(src.insulation))
        _set_cell_text(cells[7], _format_number(src.op_temp))
        _set_cell_text(cells[8], _format_number(src.op_press))


def _find_or_create_textbox(slide, box: Tuple[int, int, int, int], text: str) -> None:
  
    left, top, width, height = box
    tol = 20000 

    target = None
    for sh in slide.shapes:
        if not getattr(sh, "has_text_frame", False):
            continue
        if abs(sh.left - left) <= tol and abs(sh.top - top) <= tol and abs(sh.width - width) <= tol:
            target = sh
            break

    if target is None:
        target = slide.shapes.add_textbox(left, top, width, height)

    _set_shape_text(target, text)


def _fill_general_info_textboxes(slide, data: EquipmentData) -> None:

    _find_or_create_textbox(slide, GENERAL_DESC_BOX, data.description)
    _find_or_create_textbox(slide, GENERAL_TAG_BOX, data.tag_no)
    _find_or_create_textbox(slide, GENERAL_PMT_BOX, data.pmt_no)


def _delete_shape(shape) -> None:
    el = shape._element
    el.getparent().remove(el)


def _replace_equipment_picture(slide, image_path: Path) -> None:

    if not image_path.exists():
        return

    pictures: List[Picture] = [sh for sh in slide.shapes if isinstance(sh, Picture)]
    if not pictures:
        return

   
    pics_sorted = sorted(pictures, key=lambda p: int(p.width) * int(p.height), reverse=True)
    base = pics_sorted[0]
    left, top, width, height = base.left, base.top, base.width, base.height

    
    base_area = int(width) * int(height)
    for pic in pics_sorted:
        area = int(pic.width) * int(pic.height)
        if area >= int(base_area * 0.8):
            _delete_shape(pic)

    slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)


def _ensure_ppt_exists(ppt_abs_path: Path) -> None:
    ppt_abs_path.parent.mkdir(parents=True, exist_ok=True)
    if ppt_abs_path.exists():
        return
    if not INSPECTION_TEMPLATE_PATH.exists():
        raise FileNotFoundError(f"PowerPoint template not found at {INSPECTION_TEMPLATE_PATH}")
    shutil.copy2(INSPECTION_TEMPLATE_PATH, ppt_abs_path)



def sync_all_slides_from_masterfile(
    pptx_rel_path: str,
    workbook_rel_path: str,
    image_map: Optional[Dict[Tuple[str, str], str]] = None,
) -> Path:
 
    media_root = Path(settings.MEDIA_ROOT)

    ppt_abs = media_root / pptx_rel_path
    _ensure_ppt_exists(ppt_abs)

    prs = Presentation(str(ppt_abs))
    masterfile_path = media_root / workbook_rel_path

    for (pmt_no, eq_no), slide_idx in _EQUIPMENT_SLIDE_MAP.items():
        if slide_idx >= len(prs.slides):
            continue

        try:
            equipment_data = _load_equipment_data_from_masterfile(masterfile_path, eq_no, pmt_no)
        except Exception as e:
            print(f"[PPT Sync] Skip {pmt_no} / {eq_no}: {e}")
            continue

        slide = prs.slides[slide_idx]

      
        _fill_general_info_textboxes(slide, equipment_data)

        table = _find_material_table(slide)
        _fill_material_table(table, equipment_data)

        
        if image_map:
            rel_img = image_map.get((pmt_no, eq_no))
            if rel_img:
                _replace_equipment_picture(slide, media_root / rel_img)

    prs.save(str(ppt_abs))
    return ppt_abs


def build_inspection_plan_pptx(
    pptx_rel_path: str,
    workbook_rel_path: str,
    original_filename: str,
    slide_image_paths: Optional[List[str]] = None,
) -> Path:
   
    pmt_no, eq_no = _parse_filename(original_filename)

    image_map = None
    if slide_image_paths:
        image_map = {(pmt_no, eq_no): slide_image_paths[0]}

    return sync_all_slides_from_masterfile(
        pptx_rel_path=pptx_rel_path,
        workbook_rel_path=workbook_rel_path,
        image_map=image_map,
    )
